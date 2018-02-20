#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Created on Fri Jul  7 12:10:30 2017

@author: peter
"""

# Imports
import os # file i/o operations
import glob # file i/o operations
from datetime import date# finding today's date
from pptx import Presentation # creating the PPT
from pptx.util import Inches
import comtypes.client # for opening PowerPoint from python
import sys

""" Assign variable names to arguments from MATLAB  """
path_images = sys.argv[1]
path_reports = sys.argv[2]
batch_name = sys.argv[3]

""" Use lines 26-28 for debugging """
#path_images = 'D:\Data_Matlab\Batch_images'
#path_reports = 'D:\Data_Matlab\Reports'
#batch_name = 'batch6'

def PPTtoPDF(inputFileName, outputFileName, formatType = 32):
    """
    Converts a PPT file to a PDF by opening PowerPoint, opening the file, and 
    then saving as a PowerPoint. Requires Windows. 
    """
    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
    powerpoint.Visible = 1

    if outputFileName[-3:] != 'pdf':
        outputFileName = outputFileName + ".pdf"
    deck = powerpoint.Presentations.Open(inputFileName)
    deck.SaveAs(outputFileName, formatType) # formatType = 32 for ppt to pdf
    deck.Close()
    powerpoint.Quit()

def addImageSlide(image_file_name):
    """
    Adds a full-screen image to a blank full-screen slide.
    """
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    slide.shapes.add_picture(image_file_name, 0, 0, height=prs.slide_height, width=prs.slide_width)

# Get today's date, formatted to MATLAB's default (e.g. 2017-Jul-09)
today = date.today().strftime('%d-%b-%Y')

# make filename
reportFile = today + '_report.pptx'

# Initialize presentation
prs = Presentation()
prs.slide_height = 5143500 # Widescreen aspect ratio
title_slide_layout = prs.slide_layouts[0] # add title slide
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

# Create title slide
title.text = "Current Cycling Progress"
subtitle.text = today

# CD to directory with most recent images
os.chdir(path_images + '\\' + batch_name + '\\')

# Add .png files in this directory. Start with summary figures
all_images = glob.glob('*.png')
for file in all_images:
    if "summary" in file:
        addImageSlide(file)
        
# Cell "spec sheets"
for file in all_images:
    if "summary" not in file:
        addImageSlide(file)

# Change to directory for saving reports
os.chdir(path_reports)

# Create file names
reportFileFull = path_reports + '\\' + reportFile
reportFileFullPDF = path_reports + '\\' + reportFile.replace('pptx','pdf')

# Save powerpoint
prs.save(reportFileFull)
# Convert to PDF
PPTtoPDF(reportFileFull,reportFileFullPDF)